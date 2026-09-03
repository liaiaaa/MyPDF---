# -*- coding: utf-8 -*-
"""
============================================================
 我的专属安全 PDF 工具箱 —— 核心功能模块 (core.py)
============================================================
 纯本地 PDF 处理引擎，不依赖任何 GUI 组件。
 与 ui.py（界面）、pdf_toolbox.py（入口）共同组成完整工具箱。

 安全承诺：
   · 全部功能 100% 本地处理，代码不联网、不上传任何文件
   · 不打包、不加壳、无数字签名伪造 —— 源码完全可见、可审计
   · 仅依赖知名开源库：PyMuPDF / python-docx / openpyxl /
     python-pptx / Pillow(pptx内部用)
   · 不写入注册表、不创建计划任务、不常驻后台
============================================================
"""

import os
import re
import json
import sys
import io
import tempfile

# PyMuPDF：PDF 解析/渲染/合并/拆分/加密 + 表格提取
try:
    import pymupdf as fitz
except ImportError:
    import fitz

# 图片处理：优先用 PyMuPDF（体积小），python-pptx 需要 Pillow 时再引入
try:
    from PIL import Image
    _HAS_PIL = True
except Exception:
    _HAS_PIL = False

try:
    from docx import Document as DocxDocument
    _HAS_DOCX = True
except Exception:
    _HAS_DOCX = False

try:
    from openpyxl import Workbook
    _HAS_XLSX = True
except Exception:
    _HAS_XLSX = False

try:
    from pptx import Presentation
    _HAS_PPTX = True
except Exception:
    _HAS_PPTX = False


APP_TITLE = "我的专属安全 PDF 工具箱"
APP_VER = "1.3"
CONFIG_FILE = "_config.json"


# ============================================================
#  一、工具函数
# ============================================================

def _app_dir():
    """程序所在目录：exe 运行时取 exe 目录，源码运行时取当前目录"""
    if getattr(sys, "frozen", False):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def safe_stem(path):
    """返回不带扩展名的文件名，兼容中文路径"""
    return os.path.splitext(os.path.basename(path))[0]


def auto_out_path(src, out_dir, suffix, ext):
    """自动生成输出文件路径（避免覆盖）"""
    base = safe_stem(src)
    p = os.path.join(out_dir, f"{base}{suffix}{ext}")
    n = 1
    while os.path.exists(p):
        p = os.path.join(out_dir, f"{base}{suffix}_{n}{ext}")
        n += 1
    return p


def ensure_dir(d):
    os.makedirs(d, exist_ok=True)
    return d


def fmt_size(n):
    for u in ("B", "KB", "MB", "GB"):
        if n < 1024:
            return f"{n:.1f} {u}"
        n /= 1024
    return f"{n:.1f} TB"


def compress_quality_to_params(level):
    """压缩级别 -> (目标图片最长边像素, JPEG质量)"""
    return {"低": (3000, 80), "中": (1800, 65), "高": (1000, 50)}.get(level, (1800, 65))


def _clean_table(t):
    """去掉全空行"""
    return [row for row in t if any((c or "").strip() for c in row)]


def _page_tables(page):
    """提取页面表格：先按线条策略识别；无边框表格回退到文本对齐策略"""
    attempts = [
        {},  # 默认（lines 策略，识别有边框表格）
        {"vertical_strategy": "text", "horizontal_strategy": "text"},  # 无边框表格
    ]
    for kwargs in attempts:
        try:
            tabs = page.find_tables(**kwargs)
            tables = [_clean_table(t.extract()) for t in tabs.tables]
            # 只保留像表格的结构：至少 2 行且至少 2 列
            tables = [t for t in tables
                      if len(t) >= 2 and max((len(r) for r in t), default=0) >= 2]
            if tables:
                return tables
        except Exception:
            continue
    return []


def _image_size(path):
    """用 PyMuPDF 读取图片真实像素尺寸（替代 Pillow）"""
    try:
        pix = fitz.Pixmap(path)
        w, h = pix.width, pix.height
        pix = None
        return w, h
    except Exception:
        try:
            img = Image.open(path)
            w, h = img.size
            img.close()
            return w, h
        except Exception:
            return 0, 0


def load_config():
    try:
        p = os.path.join(_app_dir(), CONFIG_FILE)
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}


def save_config(cfg):
    try:
        p = os.path.join(_app_dir(), CONFIG_FILE)
        with open(p, "w", encoding="utf-8") as f:
            json.dump(cfg, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


# ============================================================
#  二、核心功能（纯本地）
# ============================================================

def pdf_info(path):
    """PDF 信息：页数 / 大小 / 加密状态 / 元数据"""
    doc = fitz.open(path)
    try:
        info = {
            "文件": path,
            "页数": doc.page_count,
            "文件大小": fmt_size(os.path.getsize(path)),
            "已加密": "是" if doc.is_encrypted else "否",
            "标题": doc.metadata.get("title") or "（无）",
            "作者": doc.metadata.get("author") or "（无）",
            "创建时间": doc.metadata.get("creationDate") or "（无）",
        }
        if doc.is_encrypted:
            info["提示"] = "该 PDF 已加密，使用前需先解锁。"
        return info
    finally:
        doc.close()


def pdf_merge(paths, out_path):
    """合并多个 PDF"""
    merged = fitz.open()
    try:
        for p in paths:
            with fitz.open(p) as d:
                merged.insert_pdf(d)
        merged.save(out_path, garbage=4, deflate=True)
    finally:
        merged.close()
    return out_path


def pdf_split(path, out_dir, pages=None):
    """拆分：pages=None 每页单独一个；否则页码范围 [a,b] 一个文件"""
    src = fitz.open(path)
    try:
        n = src.page_count
        if pages is None:
            outs = []
            for i in range(n):
                d = fitz.open()
                d.insert_pdf(src, from_page=i, to_page=i)
                p = os.path.join(out_dir, f"{safe_stem(path)}_第{i+1:02d}页.pdf")
                d.save(p, garbage=4, deflate=True)
                d.close()
                outs.append(p)
            return outs
        else:
            a, b = pages
            b = min(b, n)
            if a < 1 or a > b:
                raise ValueError("页码范围无效")
            d = fitz.open()
            d.insert_pdf(src, from_page=a-1, to_page=b-1)
            p = os.path.join(out_dir, f"{safe_stem(path)}_第{a}-{b}页.pdf")
            d.save(p, garbage=4, deflate=True)
            d.close()
            return [p]
    finally:
        src.close()


def _rgb_from_pixmap(pix):
    """把任意颜色空间/带 alpha 的 Pixmap 统一转为纯 RGB"""
    if pix.colorspace is not None and pix.colorspace.n != 3:
        pix = fitz.Pixmap(fitz.csRGB, pix)
    if pix.alpha:
        pix = fitz.Pixmap(pix, 0)  # 去掉 alpha（JPEG 不支持）
    return pix


def _replace_image(src, xref, pil_img, quality):
    """把 PIL 图像以 JPEG 替换回 PDF，并同步尺寸/过滤器/颜色空间"""
    buf = io.BytesIO()
    pil_img.save(buf, "JPEG", quality=quality, optimize=True)
    # compress=False 存原始 JPEG（默认会把流再 Flate 压缩，与 DCTDecode 冲突）
    src.update_stream(xref, buf.getvalue(), new=False, compress=False)
    src.xref_set_key(xref, "Filter", "/DCTDecode")      # 数据现在是 JPEG
    src.xref_set_key(xref, "DecodeParms", "null")        # 清除旧的 JPEG/PNG 参数
    src.xref_set_key(xref, "Width", str(pil_img.width))
    src.xref_set_key(xref, "Height", str(pil_img.height))
    src.xref_set_key(xref, "ColorSpace", "/DeviceRGB")   # 统一转 RGB 后必须改
    src.xref_set_key(xref, "SMask", "null")              # JPEG 无透明，去掉残留遮罩
    src.xref_set_key(xref, "Decode", "null")             # CMYK 图残留的反相数组，RGB 下必须清掉


def pdf_compress(path, out_path, level="中"):
    """压缩：图片任意比例降采样 + JPEG 二次压缩，可明显减小体积"""
    src = fitz.open(path)
    try:
        target, quality = compress_quality_to_params(level)
        for pno in range(src.page_count):
            page = src[pno]
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    pix = fitz.Pixmap(src, xref)
                except Exception:
                    continue
                pix = _rgb_from_pixmap(pix)
                # 计算目标尺寸（任意比例）
                nw, nh = pix.width, pix.height
                if max(nw, nh) > target:
                    scale = target / max(nw, nh)
                    nw = max(1, int(nw * scale))
                    nh = max(1, int(nh * scale))
                # 主路径：PIL 任意比例缩放（PIL 是 python-pptx 固定依赖，始终可用）
                try:
                    pil_img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                    if (nw, nh) != (pix.width, pix.height):
                        pil_img = pil_img.resize((nw, nh), Image.LANCZOS)
                    _replace_image(src, xref, pil_img, quality)
                except Exception:
                    # 降级：PyMuPDF 原生整数倍缩放 + JPEG 重编码
                    factor = max(2, pix.width // nw, pix.height // nh)
                    if factor > 1:
                        pix.shrink(factor)
                    src.update_stream(xref, pix.tobytes("jpeg", quality),
                                      new=False, compress=False)
                    src.xref_set_key(xref, "Filter", "/DCTDecode")
                    src.xref_set_key(xref, "DecodeParms", "null")
                    src.xref_set_key(xref, "Width", str(pix.width))
                    src.xref_set_key(xref, "Height", str(pix.height))
                    src.xref_set_key(xref, "ColorSpace", "/DeviceRGB")
                    src.xref_set_key(xref, "Decode", "null")
                pix = None
        src.save(out_path, garbage=4, deflate=True)
        return out_path
    finally:
        src.close()


def pdf_encrypt(path, out_path, password):
    """加密（AES-256）"""
    src = fitz.open(path)
    try:
        src.save(out_path, encryption=fitz.PDF_ENCRYPT_AES_256,
                 user_pw=password, owner_pw=password, permissions=0)
        return out_path
    finally:
        src.close()


def pdf_decrypt(path, out_path, password=""):
    """解密（移除密码）；已加密且无密码时自动尝试空密码"""
    src = fitz.open(path)
    try:
        if src.is_encrypted:
            if not src.authenticate(password):
                raise ValueError("密码错误，无法解锁该 PDF")
        src.save(out_path, garbage=4, deflate=True)
        return out_path
    finally:
        src.close()


def pdf_to_images(path, out_dir, dpi=150, fmt="png"):
    """PDF 转图片（每页一张）"""
    src = fitz.open(path)
    try:
        outs = []
        for i, page in enumerate(src):
            pix = page.get_pixmap(dpi=dpi)
            p = os.path.join(out_dir, f"{safe_stem(path)}_第{i+1:02d}页.{fmt}")
            pix.save(p)
            outs.append(p)
        return outs
    finally:
        src.close()


def images_to_pdf(image_paths, out_path):
    """图片转 PDF（一张图一页）"""
    doc = fitz.open()
    try:
        for ip in image_paths:
            w, h = _image_size(ip)
            if w <= 0 or h <= 0:
                raise ValueError(f"无法读取图片：{ip}")
            page = doc.new_page(width=w, height=h)
            page.insert_image(fitz.Rect(0, 0, w, h), filename=ip)
        doc.save(out_path, garbage=4, deflate=True)
        return out_path
    finally:
        doc.close()


def pdf_to_word(path, out_path):
    """PDF 转 Word：文本 + 基础表格 + 页面图片（尽量不丢内容）"""
    if not _HAS_DOCX:
        raise RuntimeError("缺少 python-docx 库，无法转 Word")
    document = DocxDocument()
    src = fitz.open(path)
    try:
        for pno, page in enumerate(src):
            # 含内嵌图片的页：先把整页渲染成图插入，保证原图不丢
            if page.get_images(full=True):
                pix = page.get_pixmap(dpi=110)
                document.add_picture(io.BytesIO(pix.tobytes("png")))
                document.add_paragraph()
            # 文本（按块取，尽量保持行序）
            blocks = page.get_text("blocks")
            for b in blocks:
                txt = (b[4] or "").strip()
                if not txt:
                    continue
                for para in txt.split("\n"):
                    if para.strip():
                        document.add_paragraph(para.strip())
            # 表格
            for t in _page_tables(page):
                if not t:
                    continue
                tb = document.add_table(rows=len(t), cols=max(len(r) for r in t))
                tb.style = "Table Grid"
                for ri, row in enumerate(t):
                    for ci, cell in enumerate(row):
                        tb.cell(ri, ci).text = (cell or "").strip()
                document.add_paragraph()
            if pno != src.page_count - 1:
                document.add_page_break()
        document.save(out_path)
        return out_path
    finally:
        src.close()


def pdf_to_excel(path, out_path):
    """PDF 转 Excel：每页表格写入独立工作表（PyMuPDF 提取）"""
    if not _HAS_XLSX:
        raise RuntimeError("缺少 openpyxl 库，无法转 Excel")
    wb = Workbook()
    wb.remove(wb.active)
    src = fitz.open(path)
    try:
        table_count = 0
        for pno, page in enumerate(src):
            tables = _page_tables(page)
            if tables:
                for t in tables:
                    if not t:
                        continue
                    table_count += 1
                    ws = wb.create_sheet(title=f"第{pno+1}页_表{table_count}"[:31])
                    for ri, row in enumerate(t):
                        for ci, cell in enumerate(row):
                            ws.cell(row=ri+1, column=ci+1, value=(cell or "").strip())
            else:
                text = page.get_text() or ""
                ws = wb.create_sheet(title=f"第{pno+1}页"[:31])
                for li, line in enumerate(text.split("\n")):
                    ws.cell(row=li+1, column=1, value=line.strip())
        if not wb.sheetnames:
            wb.create_sheet(title="空")
        wb.save(out_path)
        return out_path
    finally:
        src.close()


def pdf_to_ppt(path, out_path, dpi=150):
    """PDF 转 PPT：每页渲染为图片放入幻灯片（统一画布尺寸）"""
    if not _HAS_PPTX:
        raise RuntimeError("缺少 python-pptx 库，无法转 PPT")
    prs = Presentation()
    src = fitz.open(path)
    tmp_files = []
    try:
        if src.page_count == 0:
            raise ValueError("PDF 没有页面，无法转换")
        # 用第一页确定整份演示文稿的统一画布尺寸
        first = src[0].get_pixmap(dpi=dpi)
        prs.slide_width = int(first.width / 96 * 914400)
        prs.slide_height = int(first.height / 96 * 914400)
        for i, page in enumerate(src):
            pix = page.get_pixmap(dpi=dpi)
            png = os.path.join(tempfile.gettempdir(),
                               f"__mypdf_ppt_{i}_{os.getpid()}.png")
            pix.save(png)
            tmp_files.append(png)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(png, 0, 0,
                                     width=int(pix.width / 96 * 914400),
                                     height=int(pix.height / 96 * 914400))
        prs.save(out_path)
        return out_path
    finally:
        for p in tmp_files:
            try:
                os.remove(p)
            except OSError:
                pass
        src.close()


def pdf_extract_images(path, out_dir):
    """提取 PDF 中嵌入的所有图片"""
    src = fitz.open(path)
    try:
        outs = []
        seen = set()
        for pno, page in enumerate(src):
            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen:
                    continue
                seen.add(xref)
                try:
                    info = src.extract_image(xref)
                except Exception:
                    continue
                ext = info["ext"]
                p = os.path.join(out_dir, f"{safe_stem(path)}_图{xref}.{ext}")
                with open(p, "wb") as f:
                    f.write(info["image"])
                outs.append(p)
        return outs
    finally:
        src.close()
