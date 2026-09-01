# -*- coding: utf-8 -*-
"""
============================================================
 我的专属安全 PDF 工具箱  (MySafePDFToolbox)  v1.1
============================================================
 作者：你本人专属定制（基于 100% 开源库构建）

 安全承诺：
   · 全部功能 100% 本地处理，代码不联网、不上传任何文件
   · 不打包、不加壳、无数字签名伪造 —— 源码完全可见、可审计
   · 仅依赖知名开源库：PyMuPDF / python-docx / openpyxl /
     python-pptx / pdfplumber / Pillow / tkinterdnd2
   · 不写入注册表、不创建计划任务、不常驻后台

 v1.1 更新（点击体验优化）：
   · 支持把 PDF / 图片直接拖进窗口，自动填入当前功能
   · 支持把文件直接拖到「启动快捷方式」上打开
   · 执行时显示进度条与状态，完成后可一键打开输出目录/文件
   · 自动记住上次使用的目录
   · 多文件功能（合并/图片转PDF）改为直观的文件列表管理
   · 输出目录智能默认到源文件所在文件夹

 运行方式：
   双击「我的PDF工具箱.lnk」（无黑框），或「启动我的PDF工具箱.bat」。
   也可命令行：python pdf_toolbox.py [可选的pdf文件]
============================================================
"""

import os
import re
import json
import sys
import threading
import traceback
from datetime import datetime

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext

# 拖拽支持（缺库时自动降级，不影响其它功能）
try:
    from tkinterdnd2 import TkinterDnD, DND_FILES
    _HAS_DND = True
except Exception:
    _HAS_DND = False

# ---------- 核心功能库（全部为开源、本地库） ----------
try:
    import pymupdf as fitz   # PyMuPDF 新接口（PDF 解析/渲染/合并/拆分/加密）
except ImportError:
    import fitz
from PIL import Image  # Pillow：图片处理

try:
    import pdfplumber  # 表格提取（PDF 转 Excel）
    _HAS_PLUMBER = True
except Exception:
    _HAS_PLUMBER = False

try:
    from docx import Document as DocxDocument
    _HAS_DOCX = True
except Exception:
    _HAS_DOCX = False

try:
    from openpyxl import Workbook
    _HAS_XLSX = True
except Exception:
    _HAS_XLSX = False

try:
    from pptx import Presentation
    _HAS_PPTX = True
except Exception:
    _HAS_PPTX = False

APP_TITLE = "我的专属安全 PDF 工具箱"
APP_VER = "1.1"
CONFIG_FILE = "_config.json"


# ============================================================
#  一、工具函数
# ============================================================

def safe_stem(path):
    """返回不带扩展名的文件名，兼容中文路径"""
    return os.path.splitext(os.path.basename(path))[0]


def auto_out_path(src, out_dir, suffix, ext):
    """自动生成输出文件路径（避免覆盖）"""
    base = safe_stem(src)
    p = os.path.join(out_dir, f"{base}{suffix}{ext}")
    n = 1
    while os.path.exists(p):
        p = os.path.join(out_dir, f"{base}{suffix}_{n}{ext}")
        n += 1
    return p


def ensure_dir(d):
    os.makedirs(d, exist_ok=True)
    return d


def fmt_size(n):
    for u in ("B", "KB", "MB", "GB"):
        if n < 1024:
            return f"{n:.1f} {u}"
        n /= 1024
    return f"{n:.1f} TB"


def compress_quality_to_params(level):
    """压缩级别 -> (目标图片最长边像素, JPEG质量)"""
    return {"低": (3000, 80), "中": (1800, 65), "高": (1000, 50)}.get(level, (1800, 65))


def load_config():
    try:
        p = os.path.join(os.path.dirname(os.path.abspath(__file__)), CONFIG_FILE)
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}


def save_config(cfg):
    try:
        p = os.path.join(os.path.dirname(os.path.abspath(__file__)), CONFIG_FILE)
        with open(p, "w", encoding="utf-8") as f:
            json.dump(cfg, f, ensure_ascii=False, indent=2)
    except Exception:
        pass


# ============================================================
#  二、核心功能（纯本地）
# ============================================================

def pdf_info(path):
    """PDF 信息：页数 / 大小 / 加密状态 / 元数据"""
    doc = fitz.open(path)
    try:
        info = {
            "文件": path,
            "页数": doc.page_count,
            "文件大小": fmt_size(os.path.getsize(path)),
            "已加密": "是" if doc.is_encrypted else "否",
            "标题": doc.metadata.get("title") or "（无）",
            "作者": doc.metadata.get("author") or "（无）",
            "创建时间": doc.metadata.get("creationDate") or "（无）",
        }
        if doc.is_encrypted:
            info["提示"] = "该 PDF 已加密，使用前需先解锁。"
        return info
    finally:
        doc.close()


def pdf_merge(paths, out_path):
    """合并多个 PDF"""
    merged = fitz.open()
    try:
        for p in paths:
            with fitz.open(p) as d:
                merged.insert_pdf(d)
        merged.save(out_path, garbage=4, deflate=True)
    finally:
        merged.close()
    return out_path


def pdf_split(path, out_dir, pages=None):
    """拆分：pages=None 每页单独一个；否则页码范围 [a,b] 一个文件"""
    src = fitz.open(path)
    try:
        n = src.page_count
        if pages is None:
            outs = []
            for i in range(n):
                d = fitz.open()
                d.insert_pdf(src, from_page=i, to_page=i)
                p = os.path.join(out_dir, f"{safe_stem(path)}_第{i+1:02d}页.pdf")
                d.save(p, garbage=4, deflate=True)
                d.close()
                outs.append(p)
            return outs
        else:
            a, b = pages
            b = min(b, n)
            if a < 1 or a > b:
                raise ValueError("页码范围无效")
            d = fitz.open()
            d.insert_pdf(src, from_page=a-1, to_page=b-1)
            p = os.path.join(out_dir, f"{safe_stem(path)}_第{a}-{b}页.pdf")
            d.save(p, garbage=4, deflate=True)
            d.close()
            return [p]
    finally:
        src.close()


def pdf_compress(path, out_path, level="中"):
    """压缩：图片降采样 + 二次压缩，可明显减小体积"""
    src = fitz.open(path)
    try:
        target, quality = compress_quality_to_params(level)
        for pno in range(src.page_count):
            page = src[pno]
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    pix = fitz.Pixmap(src, xref)
                except Exception:
                    continue
                if pix.n - pix.alpha > 3:  # CMYK 先转 RGB
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                if pix.width > target or pix.height > target:
                    scale = target / max(pix.width, pix.height)
                    pix = pix.shrink(int(1 / scale)) if scale < 1 else pix
                data = pix.tobytes("jpeg", quality)
                src.update_stream(xref, data)
                pix = None
        src.save(out_path, garbage=4, deflate=True)
        return out_path
    finally:
        src.close()


def pdf_encrypt(path, out_path, password):
    """加密（AES-256）"""
    src = fitz.open(path)
    try:
        src.save(out_path, encryption=fitz.PDF_ENCRYPT_AES_256,
                 user_pw=password, owner_pw=password, permissions=0)
        return out_path
    finally:
        src.close()


def pdf_decrypt(path, out_path, password=""):
    """解密（移除密码）；已加密且无密码时自动尝试空密码"""
    src = fitz.open(path)
    try:
        if src.is_encrypted:
            if not src.authenticate(password):
                raise ValueError("密码错误，无法解锁该 PDF")
        src.save(out_path, garbage=4, deflate=True)
        return out_path
    finally:
        src.close()


def pdf_to_images(path, out_dir, dpi=150, fmt="png"):
    """PDF 转图片（每页一张）"""
    src = fitz.open(path)
    try:
        outs = []
        for i, page in enumerate(src):
            pix = page.get_pixmap(dpi=dpi)
            p = os.path.join(out_dir, f"{safe_stem(path)}_第{i+1:02d}页.{fmt}")
            pix.save(p)
            outs.append(p)
        return outs
    finally:
        src.close()


def images_to_pdf(image_paths, out_path):
    """图片转 PDF（一张图一页）"""
    doc = fitz.open()
    try:
        for ip in image_paths:
            img = Image.open(ip)
            if img.mode != "RGB":
                img = img.convert("RGB")
            w, h = img.size
            page = doc.new_page(width=w, height=h)
            page.insert_image(fitz.Rect(0, 0, w, h), filename=ip)
        doc.save(out_path, garbage=4, deflate=True)
        return out_path
    finally:
        doc.close()


def pdf_to_word(path, out_path):
    """PDF 转 Word：文本 + 基础表格"""
    if not _HAS_DOCX:
        raise RuntimeError("缺少 python-docx 库，无法转 Word")
    document = DocxDocument()
    with pdfplumber.open(path) as pdf:
        for pno, page in enumerate(pdf.pages):
            text = page.extract_text() or ""
            for para in text.split("\n"):
                if para.strip():
                    document.add_paragraph(para.strip())
            tables = page.extract_tables()
            for t in tables:
                if not t:
                    continue
                tb = document.add_table(rows=len(t), cols=max(len(r) for r in t))
                tb.style = "Table Grid"
                for ri, row in enumerate(t):
                    for ci, cell in enumerate(row):
                        tb.cell(ri, ci).text = (cell or "").strip()
                document.add_paragraph()
            if pno != len(pdf.pages) - 1:
                document.add_page_break()
    document.save(out_path)
    return out_path


def pdf_to_excel(path, out_path):
    """PDF 转 Excel：每页表格写入独立工作表"""
    if not _HAS_XLSX:
        raise RuntimeError("缺少 openpyxl 库，无法转 Excel")
    wb = Workbook()
    wb.remove(wb.active)
    with pdfplumber.open(path) as pdf:
        table_count = 0
        for pno, page in enumerate(pdf.pages):
            tables = page.extract_tables()
            if tables:
                for t in tables:
                    if not t:
                        continue
                    table_count += 1
                    ws = wb.create_sheet(title=f"第{pno+1}页_表{table_count}"[:31])
                    for ri, row in enumerate(t):
                        for ci, cell in enumerate(row):
                            ws.cell(row=ri+1, column=ci+1, value=(cell or "").strip())
            else:
                text = page.extract_text() or ""
                ws = wb.create_sheet(title=f"第{pno+1}页"[:31])
                for li, line in enumerate(text.split("\n")):
                    ws.cell(row=li+1, column=1, value=line.strip())
    if not wb.sheetnames:
        wb.create_sheet(title="空")
    wb.save(out_path)
    return out_path


def pdf_to_ppt(path, out_path, dpi=150):
    """PDF 转 PPT：每页渲染为图片放入幻灯片"""
    if not _HAS_PPTX:
        raise RuntimeError("缺少 python-pptx 库，无法转 PPT")
    prs = Presentation()
    src = fitz.open(path)
    tmp_files = []
    try:
        for i, page in enumerate(src):
            pix = page.get_pixmap(dpi=dpi)
            png = os.path.join(os.path.dirname(out_path), f"__tmp_page_{i}.png")
            pix.save(png)
            tmp_files.append(png)
            img = Image.open(png)
            pw, ph = img.size
            img.close()
            prs.slide_width = int(pw / 96 * 914400)
            prs.slide_height = int(ph / 96 * 914400)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(png, 0, 0,
                                     width=int(pw / 96 * 914400),
                                     height=int(ph / 96 * 914400))
        prs.save(out_path)
        return out_path
    finally:
        for p in tmp_files:
            try:
                os.remove(p)
            except OSError:
                pass
        src.close()


def pdf_extract_images(path, out_dir):
    """提取 PDF 中嵌入的所有图片"""
    src = fitz.open(path)
    try:
        outs = []
        seen = set()
        for pno, page in enumerate(src):
            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen:
                    continue
                seen.add(xref)
                try:
                    info = src.extract_image(xref)
                except Exception:
                    continue
                ext = info["ext"]
                p = os.path.join(out_dir, f"{safe_stem(path)}_图{xref}.{ext}")
                with open(p, "wb") as f:
                    f.write(info["image"])
                outs.append(p)
        return outs
    finally:
        src.close()


# ============================================================
#  三、GUI 界面（v1.1 点击体验优化版）
# ============================================================

FEATURES = [
    ("PDF 信息", "查看页数、大小、加密状态等基本信息"),
    ("PDF 合并", "把多个 PDF 合并成一个"),
    ("PDF 拆分", "每页拆成一个，或抽取指定页码范围"),
    ("PDF 压缩", "降低图片质量与体积，尽量保持清晰"),
    ("PDF 加密", "为 PDF 设置打开密码（AES-256）"),
    ("PDF 解密", "移除 PDF 的密码保护"),
    ("PDF 转图片", "把 PDF 每页渲染成 PNG/JPG 图片"),
    ("图片转 PDF", "把多张图片合成一个 PDF"),
    ("PDF 转 Word", "提取文本与基础表格到 .docx"),
    ("PDF 转 Excel", "提取表格数据到 .xlsx"),
    ("PDF 转 PPT", "每页转成图片幻灯片 .pptx"),
    ("提取图片", "导出 PDF 内嵌的所有图片"),
]


class FileList(tk.Frame):
    """可拖拽的多文件列表：添加 / 移除 / 清空 / 上移下移"""

    def __init__(self, master, filetypes_desc="PDF 文件", filetypes=(("PDF 文件", "*.pdf"), ("所有文件", "*.*")),
                 multi=True, allow_images=False, **kw):
        super().__init__(master, **kw)
        self.filetypes = filetypes
        self.allow_images = allow_images

        self.listbox = tk.Listbox(self, height=6, selectmode="extended",
                                  font=("Microsoft YaHei UI", 9))
        sb = ttk.Scrollbar(self, orient="vertical", command=self.listbox.yview)
        self.listbox.configure(yscrollcommand=sb.set)
        self.listbox.pack(side="left", fill="both", expand=True)
        sb.pack(side="left", fill="y")

        btns = ttk.Frame(self)
        btns.pack(side="left", fill="y", padx=(6, 0))
        ttk.Button(btns, text="添加", width=8, command=self.add).pack(pady=1)
        ttk.Button(btns, text="移除", width=8, command=self.remove).pack(pady=1)
        ttk.Button(btns, text="清空", width=8, command=self.clear).pack(pady=1)
        ttk.Button(btns, text="上移", width=8, command=lambda: self.move(-1)).pack(pady=1)
        ttk.Button(btns, text="下移", width=8, command=lambda: self.move(1)).pack(pady=1)

        if _HAS_DND:
            self.drop_target_register(DND_FILES)
            self.dnd_bind("<<Drop>>", self._on_drop)

    def _on_drop(self, event):
        files = self.root.tk.splitlist(event.data)
        self.add_paths(files)

    def add(self):
        fs = filedialog.askopenfilenames(filetypes=self.filetypes)
        if fs:
            self.add_paths(fs)

    def add_paths(self, paths):
        cur = set(self.files())
        for p in paths:
            p = os.path.normpath(p)
            if not os.path.exists(p):
                continue
            if self.allow_images and p.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff")):
                pass
            elif not self.allow_images and p.lower().endswith(".pdf"):
                pass
            else:
                continue
            if p not in cur:
                self.listbox.insert("end", p)
                cur.add(p)

    def remove(self):
        sel = list(self.listbox.curselection())
        for idx in reversed(sel):
            self.listbox.delete(idx)

    def clear(self):
        self.listbox.delete(0, "end")

    def move(self, delta):
        sel = list(self.listbox.curselection())
        if not sel:
            return
        items = list(self.listbox.get(0, "end"))
        for idx in sorted(sel, reverse=(delta > 0)):
            ni = idx + delta
            if 0 <= ni < len(items):
                items[idx], items[ni] = items[ni], items[idx]
        self.listbox.delete(0, "end")
        for it in items:
            self.listbox.insert("end", it)
        for idx in sel:
            self.listbox.selection_set(idx + delta)

    def files(self):
        return list(self.listbox.get(0, "end"))


class PdfToolboxApp:
    def __init__(self, root, initial_file=None):
        self.root = root
        self.cfg = load_config()
        self.last_dir = self.cfg.get("last_dir", "")
        self.drop_slots = []      # 已注册的"接收拖拽"的文件输入槽
        self.current_idx = 0
        self._busy = False

        root.title(f"{APP_TITLE} v{APP_VER}")
        root.geometry("980x640")
        root.minsize(900, 580)

        self._build_layout()
        self.select_feature(0)

        # 拖拽：整个窗口接收文件
        if _HAS_DND:
            root.drop_target_register(DND_FILES)
            root.dnd_bind("<<Drop>>", self.on_drop_anywhere)

        # 命令行传入的文件 → 直接查看信息
        if initial_file and os.path.exists(initial_file):
            self.root.after(300, lambda: self.open_with_file(initial_file))

        root.protocol("WM_DELETE_WINDOW", self.on_close)

    # ---------- 界面搭建 ----------
    def _build_layout(self):
        top = ttk.Frame(self.root, padding=(10, 8, 10, 4))
        top.pack(fill="x")
        ttk.Label(top, text=APP_TITLE, font=("Microsoft YaHei UI", 16, "bold")).pack(side="left")
        drop_hint = "支持把 PDF / 图片直接拖进本窗口"
        ttk.Label(top, text=("　" + drop_hint) if _HAS_DND else "　· 纯本地离线 · 源码可审计",
                  foreground="#0a7a3d").pack(side="left", padx=12)

        body = ttk.Frame(self.root)
        body.pack(fill="both", expand=True, padx=10, pady=(2, 4))

        # 左侧功能列表
        left = ttk.Frame(body, width=210)
        left.pack(side="left", fill="y")
        left.pack_propagate(False)
        ttk.Label(left, text="功能", font=("Microsoft YaHei UI", 10, "bold")).pack(anchor="w", padx=6, pady=(0, 4))
        self.listbox = tk.Listbox(left, activestyle="none", font=("Microsoft YaHei UI", 10),
                                  exportselection=False, bd=1, relief="solid")
        for name, _ in FEATURES:
            self.listbox.insert("end", name)
        self.listbox.pack(fill="both", expand=True, side="left")
        sb = ttk.Scrollbar(left, orient="vertical", command=self.listbox.yview)
        sb.pack(side="right", fill="y")
        self.listbox.configure(yscrollcommand=sb.set)
        self.listbox.bind("<<ListboxSelect>>",
                          lambda e: self.select_feature(self.listbox.curselection()[0]))

        # 右侧操作区
        self.right = ttk.Frame(body)
        self.right.pack(side="left", fill="both", expand=True, padx=(12, 0))

        # 底部：进度 + 状态 + 日志
        bottom = ttk.Frame(self.root, padding=(10, 2, 10, 8))
        bottom.pack(fill="both", expand=True)
        self.status_var = tk.StringVar(value="就绪")
        self.progress = ttk.Progressbar(bottom, mode="indeterminate", length=300)
        self.progress.pack(fill="x", pady=(0, 2))
        self.log = scrolledtext.ScrolledText(bottom, height=7, state="disabled",
                                             font=("Consolas", 9), bg="#f7f7f7")
        self.log.pack(fill="both", expand=True)
        ttk.Label(bottom, textvariable=self.status_var, foreground="#333",
                  font=("Microsoft YaHei UI", 9)).pack(anchor="w", pady=(2, 0))

    # ---------- 拖拽 ----------
    def on_drop_anywhere(self, event):
        files = [os.path.normpath(f) for f in self.root.tk.splitlist(event.data)]
        if not files:
            return
        # 若当前功能注册了文件槽，优先填入
        if self.drop_slots:
            slot = self.drop_slots[0]
            self.fill_slot(slot, files)
            self.log_write(f"已把 {len(files)} 个文件加入「{FEATURES[self.current_idx][0]}」")
            return
        # 否则：单选 PDF → 直接看信息
        self.open_with_file(files[0])

    def fill_slot(self, slot, files):
        widget = slot["widget"]
        mode = slot["mode"]
        if mode == "entry":
            widget.delete(0, "end")
            widget.insert(0, "|".join(files))
        elif mode == "filelist":
            widget.add_paths(files)

    def register_drop(self, widget, mode):
        """把某个文件输入注册为拖拽接收槽"""
        self.drop_slots.append({"widget": widget, "mode": mode})
        if _HAS_DND:
            try:
                widget.drop_target_register(DND_FILES)
                widget.dnd_bind("<<Drop>>", lambda e, w=widget, m=mode: self.fill_slot(
                    {"widget": w, "mode": m}, [os.path.normpath(x) for x in w.tk.splitlist(e.data)]))
            except Exception:
                pass

    # ---------- 通用 ----------
    def log_write(self, msg):
        self.log.configure(state="normal")
        self.log.insert("end", f"[{datetime.now():%H:%M:%S}] {msg}\n")
        self.log.see("end")
        self.log.configure(state="disabled")

    def set_status(self, msg):
        self.status_var.set(msg)

    def start_progress(self):
        self._busy = True
        self.set_status("处理中…")
        self.progress.start(10)

    def stop_progress(self, msg="完成"):
        self.progress.stop()
        self._busy = False
        self.set_status(msg)

    def browse_pdf(self):
        return filedialog.askopenfilename(
            initialdir=self.last_dir or None,
            filetypes=[("PDF 文件", "*.pdf"), ("所有文件", "*.*")])

    def browse_dir(self, title="选择输出目录"):
        d = filedialog.askdirectory(initialdir=self.last_dir or None, title=title)
        if d:
            self.last_dir = d
        return d

    def ask_out_dir(self, panel, label="输出目录（留空=源文件目录）"):
        row = ttk.Frame(panel)
        row.pack(fill="x", pady=3)
        ttk.Label(row, text=label, width=18).pack(side="left")
        var = tk.StringVar()
        ent = ttk.Entry(row, textvariable=var)
        ent.pack(side="left", fill="x", expand=True)
        btn = ttk.Button(row, text="浏览…", width=8,
                         command=lambda: var.set(self.browse_dir() or var.get()))
        btn.pack(side="left", padx=4)
        return var, ent

    def open_path(self, path):
        try:
            os.startfile(path)  # type: ignore
        except Exception as e:
            self.log_write(f"无法打开：{e}")

    def show_result(self, out_path, log_lines, is_multi=False):
        """执行成功后：记录日志 + 提供"打开输出目录/打开文件"按钮"""
        for line in log_lines:
            self.log_write(line)
        self.set_status(f"完成：{out_path}")
        # 动态加两个快捷按钮
        bar = ttk.Frame(self.right)
        bar.pack(fill="x", pady=(4, 0))
        ttk.Button(bar, text="📁 打开输出目录",
                   command=lambda: self.open_path(os.path.dirname(out_path))).pack(side="left", padx=(0, 6))
        if not is_multi and os.path.exists(out_path):
            ttk.Button(bar, text="📄 打开文件",
                       command=lambda: self.open_path(out_path)).pack(side="left")

    def run_task(self, fn, on_success):
        """后台执行，带进度条与完成回调"""

        def worker():
            try:
                result = fn()
                self.root.after(0, lambda: (self.stop_progress(), on_success(result)))
            except Exception as e:
                tb = traceback.format_exc()
                self.root.after(0, lambda: (self.stop_progress("出错"), self.log_write(f"❌ 出错：{e}"), self.log_write(tb)))

        self.start_progress()
        t = threading.Thread(target=worker, daemon=True)
        t.start()

    def add_run(self, panel, fn, on_success):
        ttk.Separator(panel).pack(fill="x", pady=8)
        row = ttk.Frame(panel)
        row.pack()
        btn = ttk.Button(row, text="▶ 开始执行", command=lambda: self.run_task(fn, on_success),
                         width=18)
        btn.pack()

    def add_single_pdf(self, panel, label="选择 PDF"):
        """标准单选 PDF 行（支持拖拽）"""
        row = ttk.Frame(panel)
        row.pack(fill="x", pady=3)
        ttk.Label(row, text=label, width=14).pack(side="left")
        var = tk.StringVar()
        ent = ttk.Entry(row, textvariable=var)
        ent.pack(side="left", fill="x", expand=True)
        ttk.Button(row, text="浏览…", width=8,
                   command=lambda: self.set_entry(ent, self.browse_pdf())).pack(side="left", padx=4)
        self.register_drop(ent, "entry")
        return var, ent

    def set_entry(self, ent, value):
        if value:
            ent.delete(0, "end")
            ent.insert(0, value)

    # ---------- 功能面板 ----------
    def select_feature(self, idx):
        self.current_idx = idx
        self.drop_slots = []
        for w in self.right.winfo_children():
            w.destroy()
        name, desc = FEATURES[idx]

        head = ttk.Frame(self.right)
        head.pack(fill="x", pady=(0, 6))
        ttk.Label(head, text=name, font=("Microsoft YaHei UI", 14, "bold")).pack(side="left")
        ttk.Label(head, text="　" + desc, foreground="#555").pack(side="left", padx=8)

        panel = ttk.Frame(self.right)
        panel.pack(fill="both", expand=True)
        self.panel = panel

        builder = getattr(self, f"ui_{idx}", None)
        if builder:
            builder(panel)
        else:
            ttk.Label(panel, text="该功能未实现").pack()

    # ---------- 各功能面板 ----------
    def ui_0(self, panel):  # PDF 信息
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        self.add_run(panel,
                     lambda: pdf_info(var.get()),
                     lambda info: (self.log_write("—— PDF 信息 ——"),
                                   [self.log_write(f"{k}：{v}") for k, v in info.items()]))

    def ui_1(self, panel):  # 合并
        ttk.Label(panel, text="添加要合并的 PDF（可多选 / 拖拽，按列表顺序合并）：",
                  foreground="#555").pack(anchor="w", padx=150)
        fl = FileList(panel, filetypes=[("PDF 文件", "*.pdf"), ("所有文件", "*.*")])
        fl.pack(fill="both", expand=True, padx=150)
        self.register_drop(fl, "filelist")
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: pdf_merge(fl.files(), os.path.join(
                         out_var.get() or os.path.dirname(fl.files()[0]),
                         f"合并结果_{datetime.now():%Y%m%d_%H%M%S}.pdf")),
                     lambda out: self.show_result(out, [f"已合并 {len(fl.files())} 个文件 → {out}（{fmt_size(os.path.getsize(out))}）"]))

    def ui_2(self, panel):  # 拆分
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        mode = tk.StringVar(value="每页拆分")
        ttk.Radiobutton(panel, text="每页拆成一个文件", variable=mode, value="每页拆分").pack(anchor="w", padx=150, pady=2)
        ttk.Radiobutton(panel, text="抽取指定页码范围（如：2-5）", variable=mode, value="范围").pack(anchor="w", padx=150, pady=2)
        rng = tk.StringVar()
        ent = ttk.Entry(panel, textvariable=rng, width=12)
        ent.pack(anchor="w", padx=230)
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._split_impl(var.get(), mode.get(), rng.get(), out_var.get()),
                     lambda outs: self.show_result(outs[0], [f"已生成 {len(outs)} 个文件 → {os.path.dirname(outs[0])}"], is_multi=True))

    def _split_impl(self, path, mode, rng, out_dir):
        if not path:
            raise ValueError("请先选择 PDF")
        out_dir = out_dir or os.path.dirname(path)
        ensure_dir(out_dir)
        if mode == "范围":
            m = re.fullmatch(r"\s*(\d+)\s*-\s*(\d+)\s*", rng)
            if not m:
                raise ValueError("页码范围格式应为 2-5")
            return pdf_split(path, out_dir, pages=(int(m.group(1)), int(m.group(2))))
        return pdf_split(path, out_dir)

    def ui_3(self, panel):  # 压缩
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        ttk.Label(panel, text="压缩程度（越高体积越小，画质下降越多）").pack(anchor="w", padx=150)
        lvl = tk.StringVar(value="中")
        for v in ("低", "中", "高"):
            ttk.Radiobutton(panel, text=v, variable=lvl, value=v).pack(side="left", padx=10)
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._compress_impl(var.get(), lvl.get(), out_var.get()),
                     lambda out: self.show_result(out, [f"压缩完成：{fmt_size(os.path.getsize(var.get()))} → {fmt_size(os.path.getsize(out))}"]))
        self._compress_src = var

    def _compress_impl(self, path, lvl, out_dir):
        if not path:
            raise ValueError("请先选择 PDF")
        out_dir = out_dir or os.path.dirname(path)
        ensure_dir(out_dir)
        out = auto_out_path(path, out_dir, "_压缩", ".pdf")
        before = os.path.getsize(path)
        pdf_compress(path, out, lvl)
        after = os.path.getsize(out)
        self._last_compress = (before, after)
        return out

    def ui_4(self, panel):  # 加密
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        row = ttk.Frame(panel); row.pack(fill="x", pady=3)
        ttk.Label(row, text="打开密码", width=14).pack(side="left")
        pw = tk.StringVar()
        ent = ttk.Entry(row, textvariable=pw, show="*")
        ent.pack(side="left", fill="x", expand=True)
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._enc_impl(var.get(), pw.get(), out_var.get()),
                     lambda out: self.show_result(out, [f"已加密（AES-256）→ {out}"]))

    def _enc_impl(self, path, pw, out_dir):
        if not path:
            raise ValueError("请先选择 PDF")
        if not pw:
            raise ValueError("请输入打开密码")
        out_dir = out_dir or os.path.dirname(path)
        ensure_dir(out_dir)
        out = auto_out_path(path, out_dir, "_加密", ".pdf")
        pdf_encrypt(path, out, pw)
        return out

    def ui_5(self, panel):  # 解密
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        row = ttk.Frame(panel); row.pack(fill="x", pady=3)
        ttk.Label(row, text="密码(若需要)", width=14).pack(side="left")
        pw = tk.StringVar()
        ent = ttk.Entry(row, textvariable=pw, show="*")
        ent.pack(side="left", fill="x", expand=True)
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._dec_impl(var.get(), pw.get(), out_var.get()),
                     lambda out: self.show_result(out, [f"已解密 → {out}"]))

    def _dec_impl(self, path, pw, out_dir):
        if not path:
            raise ValueError("请先选择 PDF")
        out_dir = out_dir or os.path.dirname(path)
        ensure_dir(out_dir)
        out = auto_out_path(path, out_dir, "_解密", ".pdf")
        pdf_decrypt(path, out, pw)
        return out

    def ui_6(self, panel):  # PDF 转图片
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        ttk.Label(panel, text="分辨率 DPI（150 清晰，300 印刷级）").pack(anchor="w", padx=150)
        dpi = tk.StringVar(value="150")
        ttk.Entry(panel, textvariable=dpi, width=10).pack(anchor="w", padx=150)
        fmt = tk.StringVar(value="png")
        ttk.Radiobutton(panel, text="PNG", variable=fmt, value="png").pack(side="left", padx=10)
        ttk.Radiobutton(panel, text="JPG", variable=fmt, value="jpg").pack(side="left", padx=10)
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._toimg_impl(var.get(), dpi.get(), fmt.get(), out_var.get()),
                     lambda outs: self.show_result(outs[0], [f"共导出 {len(outs)} 张图片 → {os.path.dirname(outs[0])}"], is_multi=True))

    def _toimg_impl(self, path, dpi, fmt, out_dir):
        if not path:
            raise ValueError("请先选择 PDF")
        out_dir = out_dir or os.path.dirname(path)
        ensure_dir(out_dir)
        return pdf_to_images(path, out_dir, dpi=int(dpi) or 150, fmt=fmt)

    def ui_7(self, panel):  # 图片转 PDF
        ttk.Label(panel, text="添加图片（可多选 / 拖拽，按列表顺序合成）：",
                  foreground="#555").pack(anchor="w", padx=150)
        fl = FileList(panel, allow_images=True,
                      filetypes=[("图片", "*.png *.jpg *.jpeg *.bmp *.gif *.webp *.tif *.tiff"), ("所有文件", "*.*")])
        fl.pack(fill="both", expand=True, padx=150)
        self.register_drop(fl, "filelist")
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._img2pdf_impl(fl.files(), out_var.get()),
                     lambda out: self.show_result(out, [f"已合成 {len(fl.files())} 张图片 → {out}"]))

    def _img2pdf_impl(self, paths, out_dir):
        if not paths:
            raise ValueError("请至少添加 1 张图片")
        out_dir = out_dir or os.path.dirname(paths[0])
        ensure_dir(out_dir)
        out = os.path.join(out_dir, f"图片合成_{datetime.now():%Y%m%d_%H%M%S}.pdf")
        images_to_pdf(paths, out)
        return out

    def ui_8(self, panel):  # PDF 转 Word
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._tofmt_impl(var.get(), out_var.get(), "word"),
                     lambda out: self.show_result(out, [f"转换完成 → {out}"]))

    def ui_9(self, panel):  # PDF 转 Excel
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._tofmt_impl(var.get(), out_var.get(), "excel"),
                     lambda out: self.show_result(out, [f"转换完成 → {out}"]))

    def ui_10(self, panel):  # PDF 转 PPT
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._tofmt_impl(var.get(), out_var.get(), "ppt"),
                     lambda out: self.show_result(out, [f"转换完成 → {out}"]))

    def _tofmt_impl(self, path, out_dir, kind):
        if not path:
            raise ValueError("请先选择 PDF")
        out_dir = out_dir or os.path.dirname(path)
        ensure_dir(out_dir)
        ext = {"word": ".docx", "excel": ".xlsx", "ppt": ".pptx"}[kind]
        out = auto_out_path(path, out_dir, "_转换", ext)
        if kind == "word":
            pdf_to_word(path, out)
        elif kind == "excel":
            pdf_to_excel(path, out)
        else:
            pdf_to_ppt(path, out)
        return out

    def ui_11(self, panel):  # 提取图片
        var, _ = self.add_single_pdf(panel, "选择 PDF")
        out_var, _ = self.ask_out_dir(panel)
        self.add_run(panel,
                     lambda: self._extract_impl(var.get(), out_var.get()),
                     lambda outs: self.show_result(outs[0], [f"共提取 {len(outs)} 张图片 → {os.path.dirname(outs[0])}"], is_multi=True))

    def _extract_impl(self, path, out_dir):
        if not path:
            raise ValueError("请先选择 PDF")
        out_dir = out_dir or os.path.dirname(path)
        ensure_dir(out_dir)
        return pdf_extract_images(path, out_dir)

    # ---------- 外部打开文件 ----------
    def open_with_file(self, path):
        path = os.path.normpath(path)
        if not os.path.exists(path):
            self.log_write(f"文件不存在：{path}")
            return
        if path.lower().endswith(".pdf"):
            self.listbox.selection_clear(0, "end")
            self.listbox.selection_set(0)
            self.select_feature(0)
            try:
                info = pdf_info(path)
                self.log_write("—— 已拖入 PDF 信息 ——")
                for k, v in info.items():
                    self.log_write(f"{k}：{v}")
                self.set_status("已载入 PDF，共 " + str(info["页数"]) + " 页")
            except Exception as e:
                self.log_write(f"读取失败：{e}")
        elif path.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp")):
            self.listbox.selection_clear(0, "end")
            self.listbox.selection_set(7)
            self.select_feature(7)
            fl = self._find_filelist(7)
            if fl:
                fl.add_paths([path])
                self.log_write(f"已把图片加入「图片转 PDF」：{path}")

    def _find_filelist(self, idx):
        # 查找当前面板中的 FileList 控件
        for w in self.right.winfo_children():
            for c in w.winfo_children():
                if isinstance(c, FileList):
                    return c
        return None

    def on_close(self):
        self.cfg["last_dir"] = self.last_dir
        save_config(self.cfg)
        self.root.destroy()


def main():
    """入口：支持把文件拖到快捷方式上打开（argv 传参）"""
    args = sys.argv[1:] if len(sys.argv) > 1 else []
    initial = args[0] if args else None
    try:
        if _HAS_DND:
            root = TkinterDnD.Tk()
        else:
            root = tk.Tk()
        app = PdfToolboxApp(root, initial_file=initial)
        root.mainloop()
    except Exception as e:
        tb = traceback.format_exc()
        try:
            log_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "run_log.txt")
            with open(log_path, "w", encoding="utf-8") as f:
                f.write(tb)
        except Exception:
            pass
        try:
            r = tk.Tk()
            r.withdraw()
            messagebox.showerror(
                "启动失败",
                "程序启动时发生错误：\n\n"
                f"{e}\n\n"
                "详细信息已保存到同目录的「run_log.txt」。\n"
                "常见原因：未安装依赖库。请执行：\n"
                "  pip install -r requirements.txt",
            )
            r.destroy()
        except Exception:
            print(tb)
        sys.exit(1)


if __name__ == "__main__":
    main()
